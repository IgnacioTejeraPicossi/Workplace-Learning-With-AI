#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AI Learning with AI - PowerPoint Presentasjon Generator (Enkel versjon)
For Legeforeningen i Oslo
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_title_slide(prs):
    """Lager forside med tittel"""
    slide_layout = prs.slide_layouts[0]  # Tittel-layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AI Learning with AI"
    subtitle.text = "Intelligent Læringsplattform for Utviklere\n\nPresentert for Legeforeningen i Oslo"
    
    return slide

def create_overview_slide(prs):
    """Lager oversiktsside over hele applikasjonen"""
    slide_layout = prs.slide_layouts[1]  # Tittel og innhold
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "📊 Komplett Oversikt - AI Learning with AI"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    # Hovedmoduler
    p1 = text_frame.paragraphs[0]
    p1.text = "🏠 Dashboard - Sentral oversikt over alle moduler"
    p1.font.size = Pt(20)
    p1.font.bold = True
    
    p2 = text_frame.add_paragraph()
    p2.text = "🔍 Repository Analyzer - Hovedprodukt for salg"
    p2.font.size = Pt(18)
    p2.font.bold = True
    
    p3 = text_frame.add_paragraph()
    p3.text = "📚 Learning Repository - Håndtering av analyser og moduler"
    p3.font.size = Pt(16)
    
    p4 = text_frame.add_paragraph()
    p4.text = "🏢 Enterprise Architecture - Avansert EA-modul"
    p4.font.size = Pt(16)
    
    p5 = text_frame.add_paragraph()
    p5.text = "❓ Quiz Generator - Automatisk evaluering"
    p5.font.size = Pt(16)
    
    p6 = text_frame.add_paragraph()
    p6.text = "🔄 Process Designer - Visuell prosessmodellering"
    p6.font.size = Pt(16)
    
    return slide

def create_repository_analyzer_slide(prs):
    """Lager side for Repository Analyzer - hovedproduktet"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🔍 Repository Analyzer - HOVEDPRODUKT FOR SALG"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    # Hovedfunksjoner
    p1 = text_frame.paragraphs[0]
    p1.text = "⭐ AUTOMATISK KODEANALYSE"
    p1.font.size = Pt(22)
    p1.font.bold = True
    
    p2 = text_frame.add_paragraph()
    p2.text = "• Analyserer GitHub-repositorier automatisk"
    p2.font.size = Pt(16)
    
    p3 = text_frame.add_paragraph()
    p3.text = "• Genererer intelligent dokumentasjon"
    p3.font.size = Pt(16)
    
    p4 = text_frame.add_paragraph()
    p4.text = "• Oppretter personlige læringsmoduler"
    p4.font.size = Pt(16)
    
    p5 = text_frame.add_paragraph()
    p5.text = "• Genererer quizzer basert på kodebase"
    p5.font.size = Pt(16)
    
    p6 = text_frame.add_paragraph()
    p6.text = "• AI-drevet innsikt og anbefalinger"
    p6.font.size = Pt(16)
    
    # Forretningsverdi
    p7 = text_frame.add_paragraph()
    p7.text = "\n💰 FORRETNINGSVERDI:"
    p7.font.size = Pt(18)
    p7.font.bold = True
    
    p8 = text_frame.add_paragraph()
    p8.text = "• Sparer 80% tid på dokumentasjon"
    p8.font.size = Pt(16)
    
    p9 = text_frame.add_paragraph()
    p9.text = "• Forbedrer kodekvalitet og vedlikehold"
    p9.font.size = Pt(16)
    
    p10 = text_frame.add_paragraph()
    p10.text = "• Automatiserer team-opplæring"
    p10.font.size = Pt(16)
    
    return slide

def create_learning_repository_slide(prs):
    """Lager side for Learning Repository"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "📚 Learning Repository - Kunnskapsadministrasjon"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    p1 = text_frame.paragraphs[0]
    p1.text = "🎯 HOVEDFUNKSJONER:"
    p1.font.size = Pt(20)
    p1.font.bold = True
    
    p2 = text_frame.add_paragraph()
    p2.text = "• Lagrer og organiserer kodeanalyser"
    p2.font.size = Pt(16)
    
    p3 = text_frame.add_paragraph()
    p3.text = "• Oppretter læringsmoduler automatisk"
    p3.font.size = Pt(16)
    
    p4 = text_frame.add_paragraph()
    p4.text = "• Genererer quizzer for evaluering"
    p4.font.size = Pt(16)
    
    p5 = text_frame.add_paragraph()
    p5.text = "• Sporer læringsutvikling"
    p5.font.size = Pt(16)
    
    p6 = text_frame.add_paragraph()
    p6.text = "• Eksporterer innhold i flere formater"
    p6.font.size = Pt(16)
    
    return slide

def create_enterprise_architecture_slide(prs):
    """Lager side for Enterprise Architecture"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🏢 Enterprise Architecture - Avansert Modul"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    p1 = text_frame.paragraphs[0]
    p1.text = "🏗️ KOMPLETTE EA-FUNKSJONER:"
    p1.font.size = Pt(20)
    p1.font.bold = True
    
    p2 = text_frame.add_paragraph()
    p2.text = "• Prosessmodellering med React Flow"
    p2.font.size = Pt(16)
    
    p3 = text_frame.add_paragraph()
    p3.text = "• Heatmap-visninger og analyser"
    p3.font.size = Pt(16)
    
    p4 = text_frame.add_paragraph()
    p4.text = "• Påvirkningsanalyse (BFS-algoritme)"
    p4.font.size = Pt(16)
    
    p5 = text_frame.add_paragraph()
    p5.text = "• Katalogadministrator for applikasjoner"
    p5.font.size = Pt(16)
    
    p6 = text_frame.add_paragraph()
    p6.text = "• Risiko- og modenhetsanalyse"
    p6.font.size = Pt(16)
    
    return slide

def create_quiz_generator_slide(prs):
    """Lager side for Quiz Generator"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "❓ Quiz Generator - Automatisk Evaluering"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    p1 = text_frame.paragraphs[0]
    p1.text = "🤖 AI-DREVET QUIZ-GENERERING:"
    p1.font.size = Pt(20)
    p1.font.bold = True
    
    p2 = text_frame.add_paragraph()
    p2.text = "• Genererer spørsmål basert på kodeanalyse"
    p2.font.size = Pt(16)
    
    p3 = text_frame.add_paragraph()
    p3.text = "• Tilpasser vanskelighetsgrad automatisk"
    p3.font.size = Pt(16)
    
    p4 = text_frame.add_paragraph()
    p4.text = "• Sporer resultater og fremgang"
    p4.font.size = Pt(16)
    
    p5 = text_frame.add_paragraph()
    p5.text = "• Støtter multiple choice og åpne spørsmål"
    p5.font.size = Pt(16)
    
    return slide

def create_technology_stack_slide(prs):
    """Lager side for teknologistack"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "💻 Teknologi og Arkitektur"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    p1 = text_frame.paragraphs[0]
    p1.text = "🛠️ FRONTEND:"
    p1.font.size = Pt(20)
    p1.font.bold = True
    
    p2 = text_frame.add_paragraph()
    p2.text = "• React.js med moderne hooks"
    p2.font.size = Pt(16)
    
    p3 = text_frame.add_paragraph()
    p3.text = "• Responsivt design og moderne UI"
    p3.font.size = Pt(16)
    
    p4 = text_frame.add_paragraph()
    p4.text = "• Chart.js for visualiseringer"
    p4.font.size = Pt(16)
    
    p5 = text_frame.add_paragraph()
    p5.text = "• React Flow for prosessmodellering"
    p5.font.size = Pt(16)
    
    p6 = text_frame.add_paragraph()
    p6.text = "\n⚙️ BACKEND:"
    p6.font.size = Pt(20)
    p6.font.bold = True
    
    p7 = text_frame.add_paragraph()
    p7.text = "• FastAPI (Python) - Høy ytelse"
    p7.font.size = Pt(16)
    
    p8 = text_frame.add_paragraph()
    p8.text = "• MongoDB - Fleksibel datalagring"
    p8.font.size = Pt(16)
    
    p9 = text_frame.add_paragraph()
    p9.text = "• AI/LLM-integrasjon for intelligent analyse"
    p9.font.size = Pt(16)
    
    return slide

def create_business_value_slide(prs):
    """Lager side for forretningsverdi"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "💰 Forretningsverdi og ROI"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    p1 = text_frame.paragraphs[0]
    p1.text = "📈 KONKRETE FORDELER:"
    p1.font.size = Pt(20)
    p1.font.bold = True
    
    p2 = text_frame.add_paragraph()
    p2.text = "• 80% reduksjon i dokumentasjonstid"
    p2.font.size = Pt(16)
    
    p3 = text_frame.add_paragraph()
    p3.text = "• 60% forbedring i kodekvalitet"
    p3.font.size = Pt(16)
    
    p4 = text_frame.add_paragraph()
    p4.text = "• 70% raskere team-opplæring"
    p4.font.size = Pt(16)
    
    p5 = text_frame.add_paragraph()
    p5.text = "• Automatisering av repetetive oppgaver"
    p5.font.size = Pt(16)
    
    p6 = text_frame.add_paragraph()
    p6.text = "\n🎯 MÅLGRUPPE:"
    p6.font.size = Pt(20)
    p6.font.bold = True
    
    p7 = text_frame.add_paragraph()
    p7.text = "• Utviklingsbedrifter"
    p7.font.size = Pt(16)
    
    p8 = text_frame.add_paragraph()
    p8.text = "• IT-konsulentselskaper"
    p8.font.size = Pt(16)
    
    p9 = text_frame.add_paragraph()
    p9.text = "• Store bedrifter med utviklingsteam"
    p9.font.size = Pt(16)
    
    return slide

def create_roadmap_slide(prs):
    """Lager side for roadmap"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "🗺️ Veikart og Fremtidige Funksjoner"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    p1 = text_frame.paragraphs[0]
    p1.text = "🚀 KOMMENDE UTVIKLING:"
    p1.font.size = Pt(20)
    p1.font.bold = True
    
    p2 = text_frame.add_paragraph()
    p2.text = "• Integrasjon med flere versjonskontrollsystemer"
    p2.font.size = Pt(16)
    
    p3 = text_frame.add_paragraph()
    p3.text = "• Avanserte AI-modeller for kodeanalyse"
    p3.font.size = Pt(16)
    
    p4 = text_frame.add_paragraph()
    p4.text = "• Team-samarbeid og deling"
    p4.font.size = Pt(16)
    
    p5 = text_frame.add_paragraph()
    p5.text = "• Mobile applikasjoner"
    p5.font.size = Pt(16)
    
    p6 = text_frame.add_paragraph()
    p6.text = "• API for tredjepartsintegrasjoner"
    p6.font.size = Pt(16)
    
    return slide

def create_contact_slide(prs):
    """Lager kontaktside"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "📞 Kontakt og Neste Steg"
    
    content = slide.placeholders[1]
    text_frame = content.text_frame
    
    p1 = text_frame.paragraphs[0]
    p1.text = "🎯 NESTE STEG:"
    p1.font.size = Pt(20)
    p1.font.bold = True
    
    p2 = text_frame.add_paragraph()
    p2.text = "• Demo av Repository Analyzer"
    p2.font.size = Pt(16)
    
    p3 = text_frame.add_paragraph()
    p3.text = "• Teknisk gjennomgang"
    p3.font.size = Pt(16)
    
    p4 = text_frame.add_paragraph()
    p4.text = "• Forretningsmodell og prising"
    p4.font.size = Pt(16)
    
    p5 = text_frame.add_paragraph()
    p5.text = "• Implementeringsplan"
    p5.font.size = Pt(16)
    
    p6 = text_frame.add_paragraph()
    p6.text = "\n💡 SPØRSMÅL?"
    p6.font.size = Pt(20)
    p6.font.bold = True
    
    p7 = text_frame.add_paragraph()
    p7.text = "Vi er klare til å svare på alle spørsmål!"
    p7.font.size = Pt(16)
    
    return slide

def main():
    """Hovedfunksjon for å lage presentasjonen"""
    print("🎯 Oppretter PowerPoint-presentasjon på norsk...")
    
    # Opprett ny presentasjon
    prs = Presentation()
    
    # Legg til lysbilder
    create_title_slide(prs)
    create_overview_slide(prs)
    create_repository_analyzer_slide(prs)
    create_learning_repository_slide(prs)
    create_enterprise_architecture_slide(prs)
    create_quiz_generator_slide(prs)
    create_technology_stack_slide(prs)
    create_business_value_slide(prs)
    create_roadmap_slide(prs)
    create_contact_slide(prs)
    
    # Lagre presentasjonen
    filename = "AI_Learning_with_AI_Presentasjon_Legeforeningen.pptx"
    prs.save(filename)
    
    print(f"✅ Presentasjon opprettet: {filename}")
    print(f"📊 Antall lysbilder: {len(prs.slides)}")
    print("🎉 Klar for Legeforeningen i Oslo!")

if __name__ == "__main__":
    main()
